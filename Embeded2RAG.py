#!/usr/bin/env python3
"""
Embeded2RAG.py

여러 소스(Excel, PPT, DB Table)를 JSON으로 정규화하고, Chunk 후
RAG 단위(주제)로 저장/인덱싱하는 도구.

설계 목표
- 임베딩 모델 미확정 상태를 고려해 "지연 로딩" 가능한 인터페이스 제공
- 벡터 DB 미확정 상태를 고려해 교체 가능한 VectorStore 팩토리 제공
- 주제별 RAG 이름(rag_name)으로 독립 저장소 관리
"""
from __future__ import annotations

import argparse
import json
import os
import pathlib
import sqlite3
from dataclasses import dataclass, asdict
from datetime import datetime
from typing import Any, Dict, Iterable, List, Optional

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter


# -----------------------------
# Optional dependencies
# -----------------------------
try:
    import pandas as pd  # Excel
except Exception:  # pragma: no cover
    pd = None

try:
    from pptx import Presentation  # PPTX
except Exception:  # pragma: no cover
    Presentation = None


# -----------------------------
# Data models
# -----------------------------
@dataclass
class RawRecord:
    source_type: str
    source_path: str
    title: str
    text: str
    metadata: Dict[str, Any]


@dataclass
class ChunkRecord:
    rag_name: str
    chunk_id: str
    text: str
    metadata: Dict[str, Any]


# -----------------------------
# Readers
# -----------------------------
def read_excel(path: str) -> List[RawRecord]:
    if pd is None:
        raise RuntimeError("pandas/openpyxl가 설치되어야 Excel을 읽을 수 있습니다.")

    records: List[RawRecord] = []
    xls = pd.ExcelFile(path)
    for sheet in xls.sheet_names:
        df = xls.parse(sheet)
        for idx, row in df.iterrows():
            row_dict = {k: _safe_json(v) for k, v in row.to_dict().items()}
            text = " | ".join(f"{k}: {v}" for k, v in row_dict.items())
            records.append(
                RawRecord(
                    source_type="excel",
                    source_path=path,
                    title=f"{os.path.basename(path)}::{sheet}::row_{idx}",
                    text=text,
                    metadata={"sheet": sheet, "row_index": int(idx), "row": row_dict},
                )
            )
    return records


def read_ppt(path: str) -> List[RawRecord]:
    if Presentation is None:
        raise RuntimeError("python-pptx가 설치되어야 PPT를 읽을 수 있습니다.")

    prs = Presentation(path)
    records: List[RawRecord] = []
    for i, slide in enumerate(prs.slides):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())

        text = "\n".join(t for t in texts if t)
        if text.strip():
            records.append(
                RawRecord(
                    source_type="ppt",
                    source_path=path,
                    title=f"{os.path.basename(path)}::slide_{i+1}",
                    text=text,
                    metadata={"slide_number": i + 1},
                )
            )
    return records


def read_db_table(sqlite_db_path: str, table_name: str, limit: Optional[int] = None) -> List[RawRecord]:
    # 기본 구현: sqlite. 추후 sqlalchemy engine으로 확장 가능.
    conn = sqlite3.connect(sqlite_db_path)
    conn.row_factory = sqlite3.Row
    try:
        query = f"SELECT * FROM {table_name}"
        if limit is not None:
            query += f" LIMIT {int(limit)}"
        rows = conn.execute(query).fetchall()
    finally:
        conn.close()

    records: List[RawRecord] = []
    for idx, row in enumerate(rows):
        row_dict = {k: _safe_json(v) for k, v in dict(row).items()}
        text = " | ".join(f"{k}: {v}" for k, v in row_dict.items())
        records.append(
            RawRecord(
                source_type="db_table",
                source_path=sqlite_db_path,
                title=f"{table_name}::row_{idx}",
                text=text,
                metadata={"table": table_name, "row_index": idx, "row": row_dict},
            )
        )
    return records


# -----------------------------
# Chunking
# -----------------------------
def to_documents(records: Iterable[RawRecord]) -> List[Document]:
    docs: List[Document] = []
    for r in records:
        meta = {
            "source_type": r.source_type,
            "source_path": r.source_path,
            "title": r.title,
            **r.metadata,
        }
        docs.append(Document(page_content=r.text, metadata=meta))
    return docs


def chunk_documents(
    docs: List[Document],
    chunk_size: int,
    chunk_overlap: int,
) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", " ", ""],
    )
    return splitter.split_documents(docs)


# -----------------------------
# Embedding / Vector store plug points
# -----------------------------
def build_embeddings(kind: str, model_name_or_path: str):
    """
    kind 예시
    - hf-local: 로컬 경로 모델
    - hf-hub: huggingface hub 모델

    실제 운영에서 중국 내 다운로드 가능한 모델로 교체 가능.
    """
    kind = kind.lower().strip()

    if kind in {"hf-local", "hf-hub"}:
        from langchain_community.embeddings import HuggingFaceEmbeddings

        # model_name_or_path에 hub id 또는 local path 전달
        return HuggingFaceEmbeddings(model_name=model_name_or_path)

    raise ValueError(
        f"지원하지 않는 embeddings kind: {kind}. "
        "원하는 벤더(예: bge, qwen embedding, jina 등)로 build_embeddings를 확장하세요."
    )


def build_vector_store(store_kind: str, persist_dir: str, embeddings, documents: List[Document]):
    store_kind = store_kind.lower().strip()

    if store_kind == "chroma":
        from langchain_community.vectorstores import Chroma

        return Chroma.from_documents(
            documents=documents,
            embedding=embeddings,
            persist_directory=persist_dir,
        )

    raise ValueError(
        f"지원하지 않는 vector store: {store_kind}. "
        "중국 내 사용 가능한 벡터 DB(Milvus/Weaviate/OpenSearch 등) 어댑터를 추가하세요."
    )


# -----------------------------
# Persistence
# -----------------------------
def ensure_rag_dirs(base_dir: str, rag_name: str) -> Dict[str, str]:
    rag_root = os.path.join(base_dir, rag_name)
    raw_json_dir = os.path.join(rag_root, "raw_json")
    chunk_json_dir = os.path.join(rag_root, "chunks")
    vectordb_dir = os.path.join(rag_root, "vectordb")

    for d in [rag_root, raw_json_dir, chunk_json_dir, vectordb_dir]:
        os.makedirs(d, exist_ok=True)

    return {
        "rag_root": rag_root,
        "raw_json_dir": raw_json_dir,
        "chunk_json_dir": chunk_json_dir,
        "vectordb_dir": vectordb_dir,
    }


def save_raw_records_json(records: List[RawRecord], out_file: str) -> None:
    with open(out_file, "w", encoding="utf-8") as f:
        json.dump([asdict(r) for r in records], f, ensure_ascii=False, indent=2)


def save_chunk_records_json(rag_name: str, docs: List[Document], out_file: str) -> None:
    chunks: List[ChunkRecord] = []
    for i, d in enumerate(docs):
        chunks.append(
            ChunkRecord(
                rag_name=rag_name,
                chunk_id=f"{rag_name}_chunk_{i}",
                text=d.page_content,
                metadata=d.metadata,
            )
        )

    with open(out_file, "w", encoding="utf-8") as f:
        json.dump([asdict(c) for c in chunks], f, ensure_ascii=False, indent=2)


def list_rags(base_dir: str) -> List[str]:
    p = pathlib.Path(base_dir)
    if not p.exists():
        return []
    return sorted([d.name for d in p.iterdir() if d.is_dir()])


# -----------------------------
# Utilities
# -----------------------------
def _safe_json(v: Any) -> Any:
    if hasattr(v, "item"):
        try:
            return v.item()
        except Exception:
            return str(v)
    if isinstance(v, (datetime,)):
        return v.isoformat()
    return v


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build RAG corpus from excel/ppt/db table")
    parser.add_argument("--rag-name", required=True, help="RAG 주제 이름 (예: finance_rag)")
    parser.add_argument("--base-dir", default="./rag_store", help="RAG 저장 루트 디렉토리")

    parser.add_argument("--excel", nargs="*", default=[], help="Excel 파일(.xlsx) 목록")
    parser.add_argument("--ppt", nargs="*", default=[], help="PPT 파일(.pptx) 목록")

    parser.add_argument("--sqlite-db", default=None, help="SQLite DB 파일 경로")
    parser.add_argument("--db-table", nargs="*", default=[], help="인덱싱할 테이블 목록")
    parser.add_argument("--db-limit", type=int, default=None, help="테이블당 최대 row")

    parser.add_argument("--chunk-size", type=int, default=800)
    parser.add_argument("--chunk-overlap", type=int, default=120)

    parser.add_argument("--embed-kind", default="hf-local", help="hf-local | hf-hub")
    parser.add_argument(
        "--embed-model",
        default="sentence-transformers/all-MiniLM-L6-v2",
        help="허깅페이스 모델 ID 또는 로컬 경로",
    )
    parser.add_argument("--vector-store", default="chroma", help="현재: chroma")

    parser.add_argument("--dry-run", action="store_true", help="벡터 DB 저장 없이 JSON/Chunk까지만 수행")
    parser.add_argument("--list-rags", action="store_true", help="생성된 RAG 목록 조회")
    return parser.parse_args()


def main() -> None:
    args = parse_args()

    if args.list_rags:
        names = list_rags(args.base_dir)
        print("RAG 목록:")
        for n in names:
            print(f"- {n}")
        return

    dirs = ensure_rag_dirs(args.base_dir, args.rag_name)

    all_records: List[RawRecord] = []

    for excel_path in args.excel:
        all_records.extend(read_excel(excel_path))

    for ppt_path in args.ppt:
        all_records.extend(read_ppt(ppt_path))

    if args.sqlite_db and args.db_table:
        for t in args.db_table:
            all_records.extend(read_db_table(args.sqlite_db, t, args.db_limit))

    if not all_records:
        raise RuntimeError("입력 데이터가 없습니다. --excel/--ppt/--sqlite-db --db-table 인자를 확인하세요.")

    docs = to_documents(all_records)
    chunks = chunk_documents(docs, chunk_size=args.chunk_size, chunk_overlap=args.chunk_overlap)

    ts = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    raw_path = os.path.join(dirs["raw_json_dir"], f"raw_{ts}.json")
    chunk_path = os.path.join(dirs["chunk_json_dir"], f"chunks_{ts}.json")
    save_raw_records_json(all_records, raw_path)
    save_chunk_records_json(args.rag_name, chunks, chunk_path)

    print(f"[OK] raw json saved: {raw_path}")
    print(f"[OK] chunk json saved: {chunk_path}")
    print(f"[INFO] chunks: {len(chunks)}")

    if args.dry_run:
        print("[DRY-RUN] 벡터 DB 저장을 건너뜁니다.")
        return

    embeddings = build_embeddings(args.embed_kind, args.embed_model)
    vector_store = build_vector_store(args.vector_store, dirs["vectordb_dir"], embeddings, chunks)

    # some stores expose persist()
    if hasattr(vector_store, "persist"):
        vector_store.persist()

    print(f"[OK] vector store saved in: {dirs['vectordb_dir']}")


if __name__ == "__main__":
    main()
